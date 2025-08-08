import argparse
from concurrent.futures import ThreadPoolExecutor
from dataclasses import dataclass
from pathlib import Path
from threading import Lock
from typing import List, Optional

from PIL import Image, ImageChops, ImageDraw
from pptx import Presentation
from langgraph.graph import END, StateGraph
from langgraph.pregel import Pregel


@dataclass
class SlideResult:
    index: int
    success: bool
    feedback: Optional[str] = None


def analyze_slide(state: dict) -> dict:
    """Analyze slide content type and render the original slide."""
    slide = state["slide"]
    img_dir: Path = state["img_dir"]
    img_dir.mkdir(parents=True, exist_ok=True)
    img_path = img_dir / f"orig_{state['index']}.jpg"

    text_shapes = 0
    has_picture = False

    img = Image.new("RGB", (1280, 720), color="white")
    draw = ImageDraw.Draw(img)
    y = 10
    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE (simplified)
            has_picture = True
        if shape.has_text_frame:
            text_shapes += 1
            draw.text((10, y), shape.text, fill="black")
            y += 20

    img.save(img_path, format="JPEG")
    state["orig_img"] = img_path

    if has_picture:
        state["content_type"] = "image"
    elif text_shapes > 2:
        state["content_type"] = "text"
    else:
        state["content_type"] = "complex"
    return state


def select_template(state: dict) -> dict:
    """Select a template layout based on content type."""
    template: Presentation = state["template"]
    tried = state.setdefault("tried_layouts", set())
    content = state.get("content_type", "text")

    layout_map = {"text": 0, "image": 1, "complex": 2}
    preferred = layout_map.get(content, 0)

    for idx in range(len(template.slide_layouts)):
        candidate = (preferred + idx) % len(template.slide_layouts)
        if candidate not in tried:
            state["layout_idx"] = candidate
            tried.add(candidate)
            break
    else:
        state["layout_idx"] = 0
    return state


def apply_layout(state: dict) -> dict:
    """Apply selected layout and copy text shapes."""
    template: Presentation = state["template"]
    slide = state["slide"]
    output: Presentation = state["output"]
    lock: Lock = state["lock"]
    layout_idx = state.get("layout_idx", 0)

    with lock:
        layout = template.slide_layouts[layout_idx]
        new_slide = output.slides.add_slide(layout)
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            textbox = new_slide.shapes.add_textbox(
                left=shape.left,
                top=shape.top,
                width=shape.width,
                height=shape.height,
            )
            textbox.text_frame.text = shape.text

    state["new_slide"] = new_slide
    return state


def mechanical_rules(state: dict) -> dict:
    """Placeholder for mechanical placement rules."""
    return apply_layout(state)


def ai_placement(state: dict) -> dict:
    """Placeholder for AI-driven placement."""
    return apply_layout(state)


def hybrid_placement(state: dict) -> dict:
    """Placeholder for hybrid placement."""
    return apply_layout(state)


def render_new_slide(state: dict) -> dict:
    """Render the new slide to JPEG."""
    slide = state["new_slide"]
    img_dir: Path = state["img_dir"]
    img_path = img_dir / f"new_{state['index']}.jpg"

    img = Image.new("RGB", (1280, 720), color="white")
    draw = ImageDraw.Draw(img)
    y = 10
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        draw.text((10, y), shape.text, fill="black")
        y += 20
    img.save(img_path, format="JPEG")

    state["new_img"] = img_path
    return state


def compare_images(state: dict) -> dict:
    """Compare original and templated images."""
    orig = Image.open(state["orig_img"])
    new = Image.open(state["new_img"])
    diff = ImageChops.difference(orig, new)
    state["quality_ok"] = diff.getbbox() is None
    return state


def generate_feedback(state: dict) -> dict:
    state["feedback"] = f"Slide {state['index']} failed layout {state.get('layout_idx')}"
    state["attempts"] = state.get("attempts", 0) + 1
    return state


def mark_result(state: dict) -> SlideResult:
    feedback = None if state.get("quality_ok", False) else state.get("feedback")
    return SlideResult(index=state["index"], success=state.get("quality_ok", False), feedback=feedback)


def build_graph() -> Pregel:
    graph = StateGraph(dict)
    graph.add_node("analyze", analyze_slide)
    graph.add_node("select", select_template)
    graph.add_node("mechanical", mechanical_rules)
    graph.add_node("ai_place", ai_placement)
    graph.add_node("hybrid", hybrid_placement)
    graph.add_node("render", render_new_slide)
    graph.add_node("compare", compare_images)
    graph.add_node("feedback", generate_feedback)

    graph.set_entry_point("analyze")
    graph.add_edge("analyze", "select")

    def route_content(state: dict) -> str:
        return state.get("content_type", "text")

    graph.add_conditional_edges(
        "select", route_content, {"text": "mechanical", "image": "ai_place", "complex": "hybrid"}
    )

    graph.add_edge("mechanical", "render")
    graph.add_edge("ai_place", "render")
    graph.add_edge("hybrid", "render")
    graph.add_edge("render", "compare")

    def route_quality(state: dict) -> str:
        return "ok" if state.get("quality_ok") else "fail"

    graph.add_conditional_edges(
        "compare", route_quality, {"ok": END, "fail": "feedback"}
    )

    def route_feedback(state: dict) -> str:
        return "retry" if state.get("attempts", 0) < 3 else "done"

    graph.add_conditional_edges(
        "feedback", route_feedback, {"retry": "select", "done": END}
    )

    return graph.compile()


def process_single_slide(
    graph: Pregel,
    template: Presentation,
    output: Presentation,
    lock: Lock,
    img_dir: Path,
    idx: int,
    slide,
) -> SlideResult:
    state = {
        "template": template,
        "slide": slide,
        "output": output,
        "lock": lock,
        "img_dir": img_dir,
        "index": idx,
        "attempts": 0,
    }
    final_state = graph.invoke(state)
    return mark_result(final_state)


def process_presentation(template_path: Path, target_path: Path, output_path: Path) -> List[SlideResult]:
    template = Presentation(template_path)
    target = Presentation(target_path)
    output = Presentation()

    graph = build_graph()
    img_dir = output_path.parent / "images"
    lock = Lock()

    results: List[SlideResult] = []
    with ThreadPoolExecutor() as executor:
        futures = [
            executor.submit(
                process_single_slide,
                graph,
                template,
                output,
                lock,
                img_dir,
                idx,
                slide,
            )
            for idx, slide in enumerate(target.slides, start=1)
        ]
        for fut in futures:
            results.append(fut.result())

    output.save(output_path)
    return results


def main(argv: Optional[List[str]] = None) -> None:
    parser = argparse.ArgumentParser(description="Generate a PPTX from template and target presentation")
    parser.add_argument("--template", required=True, type=Path, help="Path to template PPTX file")
    parser.add_argument("--target", required=True, type=Path, help="Path to target PPTX file")
    parser.add_argument("--output", type=Path, default=Path("output.pptx"), help="Destination PPTX file")
    args = parser.parse_args(argv)

    process_presentation(args.template, args.target, args.output)


if __name__ == "__main__":  # pragma: no cover - CLI entry point
    main()
